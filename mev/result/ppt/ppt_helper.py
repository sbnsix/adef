""" Power point presentation helper module """

from __future__ import annotations
import copy
import six
from pptx import Presentation


def __get_blank_slide_layout(pres: Presentation) -> object:
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]


def duplicate_slide(pres: Presentation, index: int) -> object:
    """
    Duplicate the slide with the given index in pres.

    Adds slide to the end of the presentation"""
    source = pres.slides[index]

    blank_slide_layout = __get_blank_slide_layout(pres)
    dest = pres.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")

    for key, value in six.iteritems(source.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            dest.rels.add_relationship(value.reltype, value._target, value.rId)

    return dest


def copy_slides(src_pres: Presentation, dst_pres: Presentation) -> object:
    """
    Duplicate the slide with the given index in pres.
    Adds slide to the end of the presentation
    Args:
        src_pres - source presentation object
        dst_pres - destination presentation object
    Returns:
        <dst_pres>  - Modified presentation with added slides
    """

    for slide in src_pres.slides:
        blank_slide_layout = __get_blank_slide_layout(dst_pres)
        dst = dst_pres.slides.add_slide(blank_slide_layout)

        for shp in slide.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            dst.shapes._spTree.insert_element_before(newel, "p:extLst")

        # for key, value in six.iteritems(slide.rels):
        #    # Make sure we don't copy a notesSlide relation as that won't exist
        #    if not "notesSlide" in value.reltype:
        #       dst.rels.add_relationship(value.reltype, value._target, value.rId)

    return dst_pres
