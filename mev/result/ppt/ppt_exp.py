""" Power point experiment presentation module"""

from __future__ import annotations
import glob

import pandas as pd
from pptx.util import Inches

from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from ppt_base import PresentationBase


class ExperimentPresentation(PresentationBase):
    """
    Class describing experiment presentation
    """

    def create_info(self, ad_data: dict, cd_data: dict) -> None:
        """
        Method creates detection algorithm information slide
        Args:
            ad_data:    - dict containing anomaly detection model data
            cd_data:    - dict containing cycle detector model data
        Returns:
            <None>
        """
        self.prs.local_slide_no = 1
        sld = self.prs.add_slide()
        self.prs.shape_pos(sld.shapes.title, 0.0, 2.17, 1.25, 9.0)
        self.prs.apply_style(
            sld.shapes.title,
            f"{self.algo_name} algorithm, Balance: {self.prob}",
            "Calibri (Headings)",
            44,
            True,
            False,
            PP_ALIGN.CENTER,
        )

        # Format slide and add information table
        tab = sld.shapes.add_table(
            10, 4, Inches(2.76), Inches(1.4), Inches(7.8), Inches(4.9)
        )
        tab = tab.table

        tab.autofit = False
        tab.columns[0].width = Inches(1.8)
        tab.columns[1].width = Inches(2)
        tab.columns[2].width = Inches(2)
        tab.columns[3].width = Inches(2)

        tab.cell(0, 0).merge(tab.cell(0, 3))

        self.prs.apply_style(
            tab.cell(0, 0),
            "Characteristics",
            "Calibri (Body)",
            24,
            True,
            False,
            PP_ALIGN.CENTER,
        )

        tab.cell(1, 0).merge(tab.cell(1, 1))

        # Labels on the left side of the table
        self.prs.apply_style(
            tab.cell(1, 0), "Metric", "Calibri (Body)", 24, True, False, PP_ALIGN.CENTER
        )

        tab.cell(1, 2).merge(tab.cell(1, 3))

        self.prs.apply_style(
            tab.cell(1, 2),
            "Description",
            "Calibri (Body)",
            24,
            True,
            False,
            PP_ALIGN.CENTER,
        )

        tab.cell(2, 0).merge(tab.cell(2, 1))
        tab.cell(2, 2).merge(tab.cell(2, 3))

        self.prs.apply_style(
            tab.cell(2, 0),
            "Training Time [sec]",
            "Calibri (Body)",
            24,
            True,
            False,
            PP_ALIGN.CENTER,
        )

        tab.cell(3, 0).merge(tab.cell(3, 1))
        tab.cell(3, 2).merge(tab.cell(3, 3))

        self.prs.apply_style(
            tab.cell(3, 0),
            "Evaluation Time [sec]",
            "Calibri (Body)",
            24,
            True,
            False,
            PP_ALIGN.CENTER,
        )

        self.prs.apply_style(
            tab.cell(4, 2), "AD", "Calibri (Body)", 24, True, False, PP_ALIGN.CENTER
        )
        self.prs.apply_style(
            tab.cell(4, 3), "CD", "Calibri (Body)", 24, True, False, PP_ALIGN.CENTER
        )

        tab.cell(4, 0).merge(tab.cell(4, 1))
        self.prs.apply_style(
            tab.cell(4, 0), "Metric", "Calibri (Body)", 24, True, False, PP_ALIGN.CENTER
        )

        self.prs.apply_style(
            tab.cell(5, 0),
            "Δ1 [sec]",
            "Calibri (Body)",
            24,
            True,
            False,
            PP_ALIGN.CENTER,
        )
        tab.cell(5, 2).merge(tab.cell(5, 3))
        tab.cell(5, 3).vertical_anchor = MSO_ANCHOR.MIDDLE
        tab.cell(5, 0).merge(tab.cell(5, 1))
        tab.cell(5, 0).vertical_anchor = MSO_ANCHOR.MIDDLE

        self.prs.apply_style(
            tab.cell(6, 0),
            "Δ2 [sec]",
            "Calibri (Body)",
            24,
            True,
            False,
            PP_ALIGN.CENTER,
        )
        tab.cell(6, 0).merge(tab.cell(6, 1))
        tab.cell(6, 0).vertical_anchor = MSO_ANCHOR.MIDDLE

        self.prs.apply_style(
            tab.cell(7, 0),
            "AUC [%]",
            "Calibri (Body)",
            24,
            True,
            False,
            PP_ALIGN.CENTER,
        )
        tab.cell(7, 0).merge(tab.cell(7, 1))
        tab.cell(7, 0).vertical_anchor = MSO_ANCHOR.MIDDLE

        self.prs.apply_style(
            tab.cell(8, 0), "EER", "Calibri (Body)", 24, True, False, PP_ALIGN.CENTER
        )
        tab.cell(8, 0).merge(tab.cell(8, 1))
        tab.cell(8, 0).vertical_anchor = MSO_ANCHOR.MIDDLE

        self.prs.apply_style(
            tab.cell(9, 0),
            "ACC [%]",
            "Calibri (Body)",
            24,
            True,
            False,
            PP_ALIGN.CENTER,
        )
        tab.cell(9, 0).merge(tab.cell(9, 1))
        tab.cell(9, 0).vertical_anchor = MSO_ANCHOR.MIDDLE

        # Fill in algorithm data
        tab_data = self.fill_info_data(ad_data, cd_data)

        for col in range(2, 4):
            # Row 4 is not used as it has the labels only not data
            for row in [2, 3] + list(range(5, 10)):
                if row < 4:
                    search_key = tab.cell(row, 0).text.capitalize()
                    if "[sec]" in search_key and search_key not in cd_data.keys():
                        search_key = search_key.replace("[sec]", "").strip()

                    val = f"{cd_data[search_key]:.2f}"
                    self.prs.apply_style(
                        tab.cell(row, col),
                        val,
                        "Calibri (Body)",
                        24,
                        True,
                        False,
                        PP_ALIGN.CENTER,
                    )
                else:
                    key = tab.cell(row, 0).text.lower()
                    key = (
                        key.replace("δ", "d")
                        .replace("[sec]", "")
                        .replace("[%]", "")
                        .strip()
                    )

                    # col 1 (ML) col 2 (cd)
                    col_map = {"2": "ad", "3": "cd"}
                    val = "0"
                    if 1 < col < 4:
                        if (
                            key in tab_data.keys()
                            and col_map[str(col)] in tab_data[key].keys()
                        ):
                            val = tab_data[key][col_map[str(col)]]

                    val = "" if val is None else val

                    self.prs.apply_style(
                        tab.cell(row, col),
                        val,
                        "Calibri (Body)",
                        24,
                        True,
                        False,
                        PP_ALIGN.CENTER,
                    )

        self.prs.add_slide_no(sld)
        self.prs.local_slide_no += 1

    def create_table_results(self, data: dict, data_s4: dict = None) -> None:
        """
        Method creates table with Cycle Detector numerical summary
        Args:
            data: input table data results from current simulation step
            data_s4: cycle detector data from AD model step
        Returns:
            <None>
        """
        sld = self.prs.add_slide()
        self.prs.shape_pos(sld.shapes.title, 0.0, 0.39, 1.18, 12.6)
        title = (
            (
                f"{self.algo_name} AD Model"
                if self.algo_name != "Threshold Model"
                else f"{self.algo_name.capitalize()}"
            )
            if data_s4 is not None
            else (
                f"{self.algo_name} CD Model"
                if self.algo_name != "CD Threshold Model"
                else f"{self.algo_name.capitalize()}"
            )
        )

        limit = self.cfg["process"]["limit"]
        limit = round(limit, 2)

        self.prs.apply_style(
            sld.shapes.title,
            f"{title}" f"({self.prs.local_slide_no - 1}/{self.prs.total_slides})",
            "Calibri (Headings)",
            44,
            True,
            False,
            PP_ALIGN.LEFT,
        )

        attack_types = {
            k: alg
            for (k, alg) in self.cfg["attack"]["types"].items()
            if alg["enabled"]
        }
        col_names = []
        col_names = list(attack_types.keys())

        row_names = ["Attack Name", "Δ1 [sec]"]

        if data_s4 is None:
            row_names.append("Δ2 [sec]")

        row_names += [
            "AUC [%]",
            "EER",
            "ACC [%]",
            "F1",
            "Precision",
            "Recall",
            "TPR[0]",
            f"TPR[{limit}%]",
        ]
        attack_info_columns = [
            "auc",
            "eer",
            "acc",
            "f1",
            "prec",
            "recall",
            "tpr0",
            "tprn",
            "d1",
        ]

        if data_s4 is None:
            attack_info_columns.append("d2")

        rows = len(row_names)
        cols = len(col_names) + 1
        # TODO: In case of full attack envelope shrink column size
        default_col_width = 1
        first_column_width = 2.1
        top = Inches(2)
        width_nominal = first_column_width + (cols - 1) * default_col_width
        width = Inches(width_nominal)
        left = Inches((self.prs.prs.slide_width_nominal - width_nominal) / 2)

        height = Inches(4)
        tab = sld.shapes.add_table(rows, cols, left, top, width, height)
        tab = tab.table

        tab.autofit = False
        tab.columns[0].width = Inches(first_column_width)

        # Setup column width
        for i in range(1, len(col_names)):
            tab.columns[i].width = Inches(default_col_width)

        # Top columns labels for attack names
        for i, row in enumerate(row_names):
            self.prs.apply_style(
                tab.cell(i, 0),
                row,
                "Calibri (Body)",
                18,
                True,
                False,
                PP_ALIGN.CENTER,
            )

        # Fill in data for table
        tab_data = {}

        # Fill in data for PPT table object
        if self.alg_sh in data.keys():
            for at_name, attack_info in data[self.alg_sh].items():
                tab_data[at_name] = {
                    "Δ1": {},
                }
                if data_s4 is None:
                    tab_data[at_name]["Δ2"] = {}

                tab_data[at_name].update(
                    {
                        "AUC": {},
                        "ACC": {},
                        "EER": {},
                        "F1": {},
                        "Precision": {},
                        "Recall": {},
                        "TPR[0]": {},
                        f"TPR[{limit}%]": {},
                    }
                )

                for k, val in {
                    k: v for (k, v) in attack_info.items() if k in attack_info_columns
                }.items():
                    if k.startswith("d"):
                        k1 = k.replace("d", "Δ")
                    elif k.startswith("prec"):
                        k1 = "Precision"
                    elif k.lower().startswith("recall"):
                        tab_data[at_name]["Recall"] = f"{val:.2f}"
                        continue
                    else:
                        k1 = k.upper()

                    if isinstance(val, float):
                        if "auc" == k1.lower() or "acc" == k1.lower():
                            tab_data[at_name][k1] = f"{val * 100:.2f}"
                        else:
                            if k.lower() == "tpr0":
                                k1 = "TPR[0]"
                            elif k.lower() == "tprn":
                                k1 = f"TPR[{limit}%]"
                            tab_data[at_name][k1] = f"{val:.2f}"
                    elif isinstance(val, int):
                        tab_data[at_name][k1] = f"{val}"
                    elif isinstance(val, list):
                        tab_data[at_name][k1] = f"[{val[0]:.0f}, {val[1]:.0f}]"
                    elif isinstance(val, dict):
                        if "Δ1" == k1:
                            tab_data[at_name][
                                k1
                            ] = f"[{min(data_s4[self.alg_sh][at_name][k].loc[:, 'min'].tolist()):.0f}"
                            tab_data[at_name][
                                k1
                            ] += f", {max(data_s4[self.alg_sh][at_name][k].loc[:, 'max'].tolist()):.0f}]"
                        else:
                            tab_data[at_name][
                                k1
                            ] = f"[{min(val.loc[:, 'min'].tolist()):.2f}"
                            tab_data[at_name][
                                k1
                            ] += f", {max(val.loc[:, 'max'].tolist()):.2f}]"
                    elif isinstance(val, pd.Series):
                        tab_data[at_name][k1] = f"{val.tolist()[0]:.2f}"
                    elif isinstance(val, pd.DataFrame):
                        self.log.debug(f"Columns: {val.columns}")
                        if k.lower() in val.columns:
                            tab_data[at_name][k1] = f"{val.loc[:, k.lower()].tolist()}"
                    else:
                        self.log.warn(f"Issue with data: {k1} => {val}")

        skip_labels = [
            "FBETA",
            # "PREC",
            # "RECALL",
        ]
        if data_s4 is None:
            skip_labels.append("Δ2")

        col = 1

        for label, values in tab_data.items():
            row_cnt = 0
            try:
                self.prs.apply_style(
                    tab.cell(row_cnt, col),
                    label.upper(),
                    "Calibri (Body)",
                    18,
                    True,
                    False,
                    PP_ALIGN.CENTER,
                )
            except Exception as ex:
                self.log.exception(ex)

            row_cnt = 1

            for k in values.keys():
                try:
                    self.prs.apply_style(
                        tab.cell(row_cnt, col),
                        f"{tab_data[label][k]}",
                        "Calibri (Body)",
                        18,
                        True,
                        False,
                        PP_ALIGN.CENTER,
                    )
                except Exception as ex:
                    self.log.exception(ex)

                row_cnt += 1
            col += 1

        self.prs.add_slide_no(sld)
        self.prs.local_slide_no += 1

    def experiment_heading(self) -> None:
        """
        Method adds full experiment heading
        Args:
            oa_images - original attack images used in the experiment
        Returns:
            <None>
        """
        # Metric information for main experiment slide
        exp_data = {
            "Metric": "Value",
            # Experiment ID
            "Id": f"{self.cfg['experiment']['id']}",
            # Experiment description
            "Descr": self.cfg["experiment"]["descr"],
            "Balance": self.prob,
            # Results location
            "Results path": self.cfg["experiment"]["path"],
            # Maximum number of samples taken from data set
            "Sample limit": f"{self.cfg['process']['cycle_len']*self.cfg['process']['samples']}",
            "# Production cycles": f"{self.cfg['process']['cycles']}/"
            f"{self.cfg['process']['detection_cycles']}",
        }
        sld = self.prs.add_slide()
        self.prs.shape_pos(sld.shapes.title, 0.0, 0.39, 1.18, 12.6)
        self.prs.apply_style(
            sld.shapes.title,
            "Experiment information",
            "Calibri (Headings)",
            44,
            True,
            False,
            PP_ALIGN.CENTER,
        )
        top = Inches(1.75)
        left = Inches(2.89)
        width = Inches(11.6)
        height = Inches(4)
        tab = sld.shapes.add_table(len(exp_data.keys()), 2, left, top, width, height)
        tab = tab.table
        tab.autofit = False
        tab.columns[0].width = Inches(1.8)

        # Left column labels
        row_cnt = 0
        for label, value in exp_data.items():
            self.prs.apply_style(
                tab.cell(row_cnt, 0),
                label,
                "Calibri (Body)",
                20,
                True,
                False,
                PP_ALIGN.CENTER,
            )

            self.prs.apply_style(
                tab.cell(row_cnt, 1),
                value,
                "Calibri (Body)",
                20,
                True,
                False,
                PP_ALIGN.CENTER,
            )

            row_cnt += 1

        self.prs.add_slide_no(sld)

    def attack_info(self, sam_images: dict, oa_images: dict) -> None:
        """
        Method adds attack information
        Args:
            sam_images - samples of original attack in single cycle
            oa_images - original attack images used in the experiment
        Returns:
            <None>
        """
        sam_imgx = {}

        for label, sam_image_list in sam_images.items():
            idx = 1
            for sam_image in sam_image_list:
                sam_imgx[f"{label}_{idx}"] = sam_image
                idx += 1

        self.prs.create_slide("Attack samples", sam_imgx, False)

        for img_key, img_file in oa_images.items():
            self.prs.create_slide(
                f"Original attack TS, Balance: {self.prob}", {img_key: img_file}, False
            )

        return self.prs

    def summary_slide(self, summary_images: list) -> None:
        """
        Method generates summary slide for all probability sets
        Args:
            summary_images: list of detector graphs to be displayed
        Returns:
            <None>
        """
        title = f"Algorithms comparison, Balance: {self.prob}"
        for summary_image in summary_images:
            sld = self.prs.add_slide()
            self.prs.shape_pos(sld.shapes.title, 0.0, 0.39, 1.18, 12.6)
            self.prs.apply_style(
                sld.shapes.title,
                title,
                "Calibri (Headings)",
                44,
                True,
                False,
                PP_ALIGN.CENTER,
            )
            shapes = sld.shapes

            w, h = self.prs.get_image_inch_size(summary_image)
            left = (self.prs.prs.slide_width.inches - w) / 2
            top = ((self.prs.prs.slide_height.inches - h) / 2) + 0.39
            shapes.add_picture(
                summary_image, Inches(left), Inches(top), Inches(w), Inches(h)
            )
            self.prs.add_slide_no(sld)

    def algorithm_slides(
        self,
        ad_data: dict,
        cd_data: dict,
        alg_sh: str,
        algo_name: str,
        ts_images: dict,
        ad_images: dict,
        ad_roc_images: dict,
        cd_images: dict,
        cd_roc_images: dict,
        auc_tau_images: dict,
        roc_cmp_images: dict,
        d1_images: dict,
        d2_images: dict,
        prec_recall_cmp_images: dict,
    ) -> None:
        """
        Method creates slide deck for given detection algorithm.
        Args:
            ad_data         - input AD data for the presentation
            cd_data         - input CD data for the presentation
            alg_sh          - algorithm shortcut name
            algo_name       - algorithm name
            ts_images       - dictionary containing attack name + attack image
            ad_images       - AD detection dict containing attack name + attack image
            ad_roc_images   - dict containing AD images for ROC for each attack
            cd_images       - Cycle detector dict containing attack name + attack image
            cd_roc_images   - dict containing CD images for ROC for each attack
            auc_tau_images  - AUC over Tau image overview
            roc_cmp_images  - ROC comparison between AD and CD results
            d1_images       - Δ1 delay images computed from CD models
            d2_images       - Δ2 delay images computed from CD models
            prec_recall_cmp_images - Precision/Recall generated images for CD model
        Returns:
            <None>
        """
        self.algo_name = algo_name
        self.alg_sh = alg_sh

        # Generate slides for given algorithm
        self.create_info(ad_data, cd_data)

        slides_count = int(len(ad_images.keys()) / self.prs.max_images_on_the_slide)
        num_slides = 1 if slides_count == 0 else slides_count

        if len(ts_images.keys()) % self.prs.max_images_on_the_slide > 0:
            num_slides += 1

        multiplayer = self.cfg["results"]["ppt"]["all_slides_count"]

        if self.cfg["results"]["ppt"]["detail"] != "full" and num_slides > 1:
            multiplayer -= 1

        self.prs.total_slides = num_slides * multiplayer + len(ts_images.keys())

        self.prs.create_slide(f"Manifold attack TS, Balance: {self.prob}", ts_images)

        title = f"AD model detection, Balance: {self.prob}"
        self.prs.create_slide(title, ad_images)
        self.create_table_results(ad_data, cd_data)

        if self.cfg["results"]["ppt"]["detail"] == "full":
            self.prs.create_slide(f"{title} ROC", ad_roc_images)

        self.prs.create_slide(f"CD model detection, Balance: {self.prob}", cd_images)
        self.create_table_results(cd_data)

        if self.cfg["results"]["ppt"]["detail"] == "full":
            self.prs.create_slide("CD ROC", cd_roc_images)

        self.prs.create_slide("AD model Δ1 delay performance", d1_images)

        # Comparison graphs between ROC AD and ROC CD
        self.prs.create_slide("CD model Δ2 delay performance", d2_images)

        self.prs.create_slide("CD model metrics performance", auc_tau_images)

        self.prs.create_slide(
            "CD model Precision/Recall comparison", prec_recall_cmp_images
        )
        self.prs.create_slide(
            f"CD model ROC comparison, Balance: {self.prob}", roc_cmp_images
        )

    def run(self,
            prob: str,
            ad_data: dict,
            cd_data: dict,
            file_name: str) -> None:
        """
        Method generates experiment slides
        Args:
            ad_data     - AD data input from experiment
            cd_data     - Cycle Detector data input from experiment
            file_name - name of the file where
                        presentation will be saved
        Returns:
            <None>
        """
        self.prob = f"{int(prob)}%"

        # 1. Create experiment information heading
        self.experiment_heading()

        # 2. Provide attack information slides
        attack_path = f"{self.path}attack/"

        # TODO: Validate image search using glob here:
        sam_images = PresentationBase.extract_figure_infos(glob.glob(f"{attack_path}a_data_*_sam.png"))
        oa_images = PresentationBase.extract_figure_info(
            glob.glob(f"{attack_path}a_data_*_oa.png")
        )
        self.attack_info(sam_images, oa_images)

        # 3. Iterate through all detection methods
        #    and list all results
        path = f"{self.path}graph/"
        for detector_name in self.detector_names:
            detector = detector_name.lower()
            # For each algorithm create a separate deck of slides from a
            # preconfigured template
            det_name = self.cfg["model"]["ad"][detector]["name"]

            path_ad = f"{self.path}detection/"

            # Find images
            ts_images = {} # PresentationBase.extract_figure_info(
            #    glob.glob(f"{path}a_data_*{detector}*_out_ts.png")
            #)

            ad_images = PresentationBase.extract_figure_info(
                glob.glob(f"{path}a_data_*{detector}*_det_ts.png")
            )

            ad_roc_images = {} # PresentationBase.extract_figure_info(
            #    glob.glob(f"{path_ad}a_data_*{detector}*_rcm.png")
            #)

            cd_images = PresentationBase.extract_figure_info(
                glob.glob(f"{path}a_data_*{detector}*_det_cd.png")
            )
            cd_roc_images = PresentationBase.extract_figure_info(
                glob.glob(f"{path}a_data_*{detector}*_m*_p*_roc.png")
            )

            auc_tau_images = PresentationBase.extract_figure_info(
                glob.glob(f"{path}a_data_*{detector}*_tau_c.png")
            )

            d1_images = PresentationBase.extract_figure_info(
                glob.glob(f"{path_ad}a_data_*{detector}*_d1.png")
            )

            d2_images = PresentationBase.extract_figure_info(
                glob.glob(f"{path}a_data_*{detector}*_d2.png")
            )

            roc_cmp_images = PresentationBase.extract_figure_info(
                glob.glob(f"{path}a_data_*{detector}*_auc_diff.png")
            )

            prec_recall_cmp_images = PresentationBase.extract_figure_info(
                glob.glob(f"{path}a_data_*{detector}*_prec_recall_diff.png")
            )

            self.algorithm_slides(
                ad_data,
                cd_data,
                detector,
                det_name,
                ts_images,
                ad_images,
                ad_roc_images,
                cd_images,
                cd_roc_images,
                auc_tau_images,
                roc_cmp_images,
                d1_images,
                d2_images,
                prec_recall_cmp_images,
            )

        summary_images = glob.glob(f"{path}*_sum*.png")

        # Summary slide here
        self.summary_slide(summary_images)

        # Saving presentation file
        self.prs.save(file_name)
        self.log.debug(f"PPT saved to: {file_name}")
