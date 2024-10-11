""" Power point presentation helper module """

from __future__ import annotations

import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# from pptx.util import Cm

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from img_helper import ImgHelper


class PptGenerator:
    """
    Presentation generator class used in results generation to PowerPoint
    """

    def __init__(self, cfg: dict, logger: object) -> None:
        """
        CTOR
        Args:
            cfg: experiment configuration object
        Returns:
            <None>
        """

        self.cfg = cfg
        self.log = logger
        self.algo_name = ""
        self.alg_sh = ""

        # Compute number of slides additional slides in case
        # of large number of attacks that would require more than one slide to
        # showcase all test results.
        self.total_slides = self.cfg["results"]["ppt"]["all_slides_count"]

        # PowerPoint main object entry
        self.prs = Presentation()

        # Default presentation slide size 16:9
        self.prs.slide_width_nominal = 13.30
        self.prs.slide_width = Inches(13.30)
        self.prs.slide_height_nominal = 7.5
        self.prs.slide_height = Inches(7.5)

        # Default slide layout - empty slide
        self.lyt = self.prs.slide_layouts[5]
        self.slide_no = 1
        self.local_slide_no = 1
        self.max_images_on_the_slide = 6
        # Image [top, left, width, height], Text label [top, left, width, height]
        self.layout_1_1 = [
            [[1.98, 1.11, 9.33, 6.39], [6.18, 0.53, 4.38, 0.4]],
        ]

        self.layout_1_2 = [
            [[1.95, 2.53, 4.38, 3], [2.22, 2.13, 4.38, 0.4]],
            [[6.43, 2.53, 4.38, 3], [6.7, 2.13, 4.38, 0.4]],
        ]

        self.layout_2_2 = [
            [[1.85, 1.21, 4.38, 3], [2.12, 0.81, 4.38, 0.4]],
            [[6.33, 1.21, 4.38, 3], [6.6, 0.81, 4.38, 0.4]],
            [[1.85, 4.29, 4.38, 3], [2.12, 3.81, 4.38, 0.4]],
            [[6.33, 4.29, 4.38, 3], [6.6, 3.81, 4.38, 0.4]],
        ]

        self.layout_1_3 = [
            [[0, 2.53, 4.38, 3], [0.27, 2.13, 4.38, 0.4]],
            [[4.48, 2.53, 4.38, 3], [4.75, 2.13, 4.38, 0.4]],
            [[8.95, 2.53, 4.38, 3], [9.22, 2.13, 4.38, 0.4]],
        ]

        self.layout_2_3 = [
            [[0, 1.4, 4.38, 3], [0.27, 1, 4.38, 0.4]],
            [[0, 4.48, 4.38, 3], [0.27, 4, 4.38, 0.4]],
            [[4.48, 1.4, 4.38, 3], [4.75, 1, 4.38, 0.4]],
            [[4.48, 4.48, 4.38, 3], [4.75, 4, 4.38, 0.4]],
            [[8.95, 1.4, 4.38, 3], [9.22, 1, 4.38, 0.4]],
            [[8.95, 4.48, 4.38, 3], [9.22, 4, 4.38, 0.4]],
        ]

        self.layout_2_3_1 = [
            [[0, 1.4, 4.38, 3], [0.27, 1, 4.38, 0.4]],
            [[0, 4.48, 4.38, 3], [0.27, 4, 4.38, 0.4]],
            [[4.48, 1.4, 4.38, 3], [4.75, 1, 4.38, 0.4]],
            [[4.48, 4.48, 4.38, 3], [4.75, 4, 4.38, 0.4]],
            [[8.95, 1.4, 4.38, 3], [9.22, 1, 4.38, 0.4]],
            [[8.95, 4.48, 4.38, 3], [9.22, 4, 4.38, 0.4]],
        ]

    def add_slide_no(self, slide) -> None:
        """
        Method adds slide number to slide object
        Args:
            slide: slide object where new textbox with slide number will be added
        Returns:
            <None>
        """
        left = Inches(12.91)
        top = Inches(7.16)
        width = Inches(0.34)
        height = Inches(0.4)
        txt_box = slide.shapes.add_textbox(left, top, width, height)
        tf = txt_box.text_frame.paragraphs[0]
        tf.text = f"{self.slide_no}"
        tf.alignment = PP_ALIGN.RIGHT
        tf.font.name = "Calibri (Body)"
        tf.font.size = Pt(14)

        self.slide_no += 1

    def shape_pos(
        self, shape: object, top: float, left: float, height: float, width: float
    ) -> None:
        """
        Method sets shape position and size
        Args:
            shape: shape object where title will be placed
            top: float,
            left: float,
            height: float,
            width: float,
        Returns:
            <None>
        """
        shape.left = Inches(left)
        shape.top = Inches(top)
        if width is not None:
            shape.width = Inches(width)
        if height is not None:
            shape.height = Inches(height)

    def apply_style(
        self,
        shape: object,
        text: str,
        font_name: str,
        font_size: int,
        bold: bool,
        italic: bool,
        h_alignment: PP_ALIGN,
        v_alignment: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    ) -> None:
        """
        Method applies given style to a shape object.
        Args:
            shape:
            text:
            font_name:
            font_size:
            bold:
            italic:
            h_alignment:
            v_alignment:

        Returns:

        """
        if not isinstance(text, str):
            shape.text = str(text)
        else:
            shape.text = text
        shape.vertical_anchor = v_alignment
        c = shape.text_frame.paragraphs[0]
        c.font.name = font_name
        c.font.size = Pt(font_size)
        c.font.bold = bold
        c.font.italic = italic
        c.alignment = h_alignment

    def slide_layout(self, images: dict) -> list:
        """
        Method optimize best slide layout
        Args:
            images - dict containing labels and corresponding of images to be displayed
        Returns:
            <list>  - number of corresponding layout coordinates
                      for each slide
        """
        layouts = []
        displayed_images = 0
        images_count = len(images)

        # Compute number of slides required to display all images
        number_of_slides = int(images_count / self.max_images_on_the_slide)

        # Append N slide 2x3 slides (Max configuration)
        for i in range(0, number_of_slides):
            layouts.append(self.layout_2_3)
            displayed_images += self.max_images_on_the_slide

        # For remaining group add last adjusting layout
        r_slides = images_count - displayed_images

        if r_slides > 0:
            remaining_layouts = [
                self.layout_1_1,
                self.layout_1_2,
                self.layout_1_3,
                self.layout_2_2,
                self.layout_2_3_1,
                self.layout_2_3,
            ]
            layouts.append(remaining_layouts[r_slides - 1])

        return layouts

    def create_slide(self,
                     title: str,
                     images: dict,
                     slide_count: bool = True) -> None:
        """
        Method creates slide with specific experiment images layout correctly on the slide.
        Args:
            title:  name of the slide
            images: dict containing corresponding list of images and its main label
            slide_count: include slide number on the slide
        Returns:
            <None>
        """

        layouts = self.slide_layout(images)
        start = 0

        for layout in layouts:
            stop = start + len(layout)
            imgs = {x: images[x] for x in list(images.keys())[start:stop]}
            start = stop

            sld = self.prs.slides.add_slide(self.lyt)
            self.shape_pos(sld.shapes.title, 0.0, 0.39, 1.18, 12.6)

            title_full = title
            if slide_count:
                title_full = (
                    f"{self.algo_name} {title} "
                    f"({self.local_slide_no - 1}/{self.total_slides})"
                )

            self.apply_style(
                sld.shapes.title,
                title_full,
                "Calibri (Headings)",
                44,
                True,
                False,
                PP_ALIGN.LEFT,
            )

            shapes = sld.shapes

            placements = layout
            l_idx = 0
            keys = []

            # TODO: Create the loop for multiple images
            for caption, image in imgs.items():
                if l_idx >= self.max_images_on_the_slide:
                    break
                keys.append(caption)

                placement = placements[l_idx]
                # Add image
                # https://python-pptx.readthedocs.io/en/latest/dev/analysis/shp-picture.html

                dpi = ImgHelper.get_image_dpi(image)
                width, height = ImgHelper.get_image_size(image)

                x1 = placement[0][0]
                x2 = placement[0][2]
                y1 = placement[0][1]
                y2 = placement[0][3]

                rect_width = ImgHelper.inches_to_pixels(
                    x2, dpi
                ) - ImgHelper.inches_to_pixels(x1, dpi)
                rect_height = ImgHelper.inches_to_pixels(
                    y2, dpi
                ) - ImgHelper.inches_to_pixels(y1, dpi)

                # If image width is very wide adjust it to the full width of the slide
                if "_oa." in image and width > rect_width and 1 == len(layouts):
                    x1 = 0
                    x2 = 13.3
                    if height < rect_height:
                        ppi = height / (height / dpi)
                        delta_height = ImgHelper.px_to_inches(
                            rect_height, ppi
                        ) - ImgHelper.px_to_inches(height, ppi)
                        y1 = y1 + delta_height
                        y2 = (
                            y1
                            + ImgHelper.px_to_inches(height, ppi)
                            - delta_height * 2.5
                        )

                picture = shapes.add_picture(
                    image,
                    Inches(x1),
                    Inches(y1),
                    Inches(x2),
                    Inches(y2),
                )

                """
                # Add image caption
                tb = shapes.add_textbox(
                    Inches(placement[1][0]),
                    Inches(placement[1][1]),
                    Inches(placement[1][2]),
                    Inches(placement[1][3]),
                )
                p = tb.text_frame.add_paragraph()

                if caption in self.cfg["attack"]["types"].keys():
                    p.text = self.cfg["attack"]["types"][caption]["descr"]
                else:
                    p.text = caption.capitalize()

                p.font.name = "Calibri (Body)"
                p.font.size = Pt(18)
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT
                p.left = Inches(placement[1][0])
                p.top = Inches(placement[1][1])
                """

                l_idx += 1

            self.add_slide_no(sld)
            self.local_slide_no += 1

    def create_slide_image(
        self, title: str, images: dict, image_size: list, slide_count: bool = True
    ) -> None:
        """
        Method creates slide with single image on the slide
        of specific size
        Args:
            title:  name of the title
            images: dict containing corresponding image and its label
            image_size: coordinates containing image width and image length
            slide_count: include slide count
        Returns:
            <None>
        """

        for caption, image in images.items():
            sld = self.prs.slides.add_slide(self.lyt)
            self.shape_pos(sld.shapes.title, 0.0, 0.39, 1.18, 12.6)

            title_full = f"{title} {caption}"
            if slide_count:
                title_full = (
                    f"{self.algo_name} {title} "
                    f"({self.local_slide_no - 1}/{self.total_slides})"
                )

            self.apply_style(
                sld.shapes.title,
                title_full,
                "Calibri (Headings)",
                40,
                True,
                False,
                PP_ALIGN.LEFT,
            )

            shapes = sld.shapes

            left = (self.prs.slide_width_nominal - image_size[1]) / 2
            top = self.prs.slide_height_nominal - image_size[0]
            picture = shapes.add_picture(
                image,
                Inches(left),
                Inches(top),
                Inches(image_size[1]),
                Inches(image_size[0]),
            )

            """
            if caption:
                # Add image caption
                tb = shapes.add_textbox(
                    Inches(),
                    Inches(),
                    Inches(),
                    Inches(),
                )
                p = tb.text_frame.add_paragraph()

                if caption in self.cfg["attack"]["types"].keys():
                    p.text = self.cfg["attack"]["types"][caption]["descr"]
                else:
                    p.text = caption.capitalize()

                p.font.name = "Calibri (Body)"
                p.font.size = Pt(18)
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT
                p.left = Inches(placement[1][0])
                p.top = Inches(placement[1][1])
                l_idx += 1
            """

            self.add_slide_no(sld)
            self.local_slide_no += 1

    def get_image_inch_size(self, image_file: str) -> (float, float):
        """
        Method returns width and height of an image
        Args:
            image_file: image file name to be checked
        Returns:
            <width, height> - height and width of an image expressed in inches
        """
        width = 0
        height = 0

        with Image.open(image_file) as im:
            dpi = round(im.info["dpi"][0])
            width, height = pd.Series(list(im.size)) / dpi

        return width, height

    def add_slide(self) -> object:
        """
        Method adds slide to the presentation
        Returns:
            <object> - object slide
        """
        sld = self.prs.slides.add_slide(self.lyt)
        return sld

    def save(self, file_name: str) -> None:
        """
        Method saves presentation
        Args:
            file_name: presentation file name
        Returns:
            <None>
        """
        self.prs.save(file_name)
