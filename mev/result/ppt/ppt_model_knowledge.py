""" Power point experiment presentation module"""

from __future__ import annotations
import glob

import pandas as pd

from pptx.enum.text import PP_ALIGN

# from pptx.enum.text import MSO_ANCHOR

from ppt_base import PresentationBase


class ModelKnowledgePresentation(PresentationBase):
    """
    Class describing experiment presentation
    """

    def head_slide(self, exp_name: str) -> None:
        """
        Method creates head summary slide
        Args:
            exp_name: experiment name
        Returns:
            <None>
        """
        ex_els = exp_name.split("_")

        self.prs.slide_no = 1
        sld = self.prs.add_slide()
        height_nominal = 3.18
        top_val = (self.prs.prs.slide_height_nominal - height_nominal) / 2
        self.prs.shape_pos(sld.shapes.title, top_val, 2.17, height_nominal, 9.0)
        self.prs.apply_style(
            sld.shapes.title,
            f"Model knowledge\nexperiment summary",
            "Calibri (Headings)",
            44,
            True,
            False,
            PP_ALIGN.CENTER,
        )
        self.prs.slide_no += 1

    def results_phase_1(self, meta_images: dict) -> None:
        """
        Phase 1 provides evaluation of anomaly detection effectiveness by algorithm
        over amount of data used in training vs. different types of attacks
        Args:
            meta_images: list of meta level images
        Returns:
            <None>
        """
        # Add images for each of the detection methods
        for images in meta_images.values():
            self.prs.create_slide_image("", images, [6.42, 9.37], False)

    def results_phase_2(self) -> None:
        """
        Returns:
            <None>
        """
        return

    def results_phase_3(self) -> None:
        """
        Returns:
            <None>
        """
        return

    def run(self, data: pd.DataFrame, file_name: str) -> None:
        """
        Method generates experiment slides
        Args:
            data     : DataFrame containing summary information
            file_name: PowerPoint presentation file name
        Returns:
            <None>
        """
        image_path = file_name[: file_name.rindex("/")]
        image_list = glob.glob(f"{image_path}/meta*.png")
        images = {}
        exp_id = self.cfg["experiment"]["id"].lower()

        # Split images per detector category
        for image in image_list:
            img = image[image.rindex("\\") + 1 :]
            img = img[len("meta") + len(exp_id) + 2 :]
            img_el = img.split("_")
            detector_name = f"Det: {img_el[0].upper()}, Scale: {img_el[1].lower()}"
            label = f"Det: {img_el[0].upper()}, Scale: {img_el[1].lower()}"
            images[detector_name] = {label: image}

        self.head_slide(self.cfg["experiment"]["id"].lower())

        # X - Amount of attacks (model knowledge), Y - AUC performance
        self.results_phase_1(images)

        # X - Amount of attacks (model knowledge), Y - Accuracy
        # self.results_phase_2()

        # X - Amount of attacks (model knowledge), Y - F1 score
        # self.results_phase_3()

        # Saving presentation file
        self.prs.save(file_name)
        self.log.debug(f"Model knowledge summary PPT saved to: {file_name}")
